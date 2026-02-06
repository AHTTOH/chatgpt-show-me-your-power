from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


OUTPUT_DIR = Path(__file__).resolve().parent.parent / "output"
OUTPUT_FILE = OUTPUT_DIR / "antonai_chatgpt_deepdive.pptx"

ACCENT = RGBColor(0x1E, 0x88, 0xE5)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDES = [
    {
        "title": "ChatGPT 사용법을 뼛속까지 (안똔AI)",
        "lines": ["Apps / 개발자설정 / MCP / 자동화 / 실전 데모"],
        "bullets": False,
    },
    {
        "title": "오늘 얻어갈 것 5가지",
        "lines": [
            "ChatGPT 핵심 화면 구조 이해",
            "앱과 모델의 역할 분리",
            "개발자 설정으로 품질·안전 강화",
            "MCP로 도구 연결하는 법",
            "실전 데모 2개 워크플로",
        ],
        "bullets": True,
    },
    {
        "title": "ChatGPT 화면 구조",
        "lines": [
            "모델: 답변을 만드는 두뇌",
            "툴: 계산/검색/실행 같은 손발",
            "앱: 캔버스·캘린더 등 연결 서비스",
            "파일: 입력/출력 자료",
            "프로젝트: 목적별 지식·맥락 저장소",
        ],
        "bullets": True,
    },
    {
        "title": "Apps: Canva 포함",
        "lines": [
            "앱이 하는 일: UI 제공, 외부 서비스 접속",
            "모델이 하는 일: 지시 해석, 콘텐츠 생성",
            "앱은 ‘창구’, 모델은 ‘엔진’",
            "연결 앱 목록은 사용 시점에만 권한 부여",
        ],
        "bullets": True,
    },
    {
        "title": "개발자 설정(Developer settings)",
        "lines": [
            "개발자 모드: 도구/프로젝트 권한 제어",
            "안전 필터: 위험한 호출 차단",
            "실험 기능: 최신 모델/툴 우선 사용",
            "로그 확인: 실패 원인 추적",
        ],
        "bullets": True,
    },
    {
        "title": "MCP 핵심 개념",
        "lines": [
            "도구 호출을 표준화하는 연결 프로토콜",
            "모델이 ‘요청’, 서버가 ‘실행’",
            "권한 최소화 + 읽기 전용부터 시작",
            "환경 분리(개발/운영)로 사고 방지",
        ],
        "bullets": True,
    },
    {
        "title": "실전 데모 1: 문서/표/리서치",
        "lines": [
            "1) 요구사항 요약 → 질문 리스트 생성",
            "2) 문서·표·리서치 병렬 수집",
            "3) 핵심 인사이트 5줄로 정리",
        ],
        "bullets": True,
    },
    {
        "title": "실전 데모 2: 자동화 시나리오",
        "lines": [
            "조건: 신규 파일 업로드 감지",
            "트리거: 요약 + 분류 + 알림",
            "출력: 슬랙 메시지 + 보고서 초안",
        ],
        "bullets": True,
    },
    {
        "title": "Codex란?",
        "lines": [
            "ChatGPT 채팅: 대화 중심",
            "Codex 작업공간: 파일/코드 중심",
            "버전 관리 + 자동 실행",
            "다단계 산출물에 최적화",
        ],
        "bullets": True,
    },
    {
        "title": "Codex로 산출물 만들기",
        "lines": [
            "코드 → 문서 → 슬라이드까지 연결",
            "목표·제약·출력물 3줄로 정의",
            "반복 피드백으로 ‘끝까지’ 완주",
            "오늘 실습: 슬라이드 자동 생성",
        ],
        "bullets": True,
    },
    {
        "title": "체크리스트",
        "lines": [
            "실패: 모호한 목적, 범위 없음",
            "실패: 출력 형식 지시 누락",
            "성공: 목적/대상/형식/제약 명시",
            "성공: 예시 제공 + 검증 기준",
            "성공: 단계별 확인 요청",
        ],
        "bullets": True,
    },
    {
        "title": "엔딩",
        "lines": [
            "다음 영상: MCP 실전 서버 만들기",
            "구독/자료 링크 자리",
            "질문은 댓글에!",
        ],
        "bullets": True,
    },
]


def build_deck() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    blank_layout = presentation.slide_layouts[6]
    title_left = Inches(0.7)
    title_top = Inches(0.3)
    title_width = Inches(12.0)
    title_height = Inches(0.6)

    line_left = Inches(0.7)
    line_top = Inches(0.98)
    line_width = Inches(12.0)
    line_height = Inches(0.04)

    body_left = Inches(0.9)
    body_top = Inches(1.3)
    body_width = Inches(11.6)
    body_height = Inches(5.6)

    for slide_data in SLIDES:
        slide = presentation.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(
            title_left, title_top, title_width, title_height
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.text = slide_data["title"]
        title_run = title_frame.paragraphs[0].runs[0]
        title_run.font.size = Pt(32)
        title_run.font.bold = True
        title_run.font.color.rgb = BLACK

        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            line_left,
            line_top,
            line_width,
            line_height,
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

        body_box = slide.shapes.add_textbox(
            body_left, body_top, body_width, body_height
        )
        body_frame = body_box.text_frame
        body_frame.clear()
        for idx, line_text in enumerate(slide_data["lines"]):
            paragraph = body_frame.paragraphs[0] if idx == 0 else body_frame.add_paragraph()
            paragraph.text = line_text
            paragraph.level = 0
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = BLACK
            paragraph.bullet = slide_data["bullets"]

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    presentation.save(OUTPUT_FILE)


if __name__ == "__main__":
    build_deck()
    print(f"Generated: {OUTPUT_FILE}")
